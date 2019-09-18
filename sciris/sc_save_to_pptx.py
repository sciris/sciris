import pptx as pp
import os

def savepptx(filename=None, template=None, slides=None, image_path=None):
    '''
    :param filename: A name for the desired output document. It should end in .pptx.
    :param template: The filepath to the powerpoint template which is to be used.
    :param slides: A list of dicts, with each dict specifying the attributes to be added in a slide. The format of each
    dict should match with the construction of an appropriate slide within your template.
    :param image_path: The filepath to any images which are to be added to a slide.

    If no template powerpoint is provided then an automatic template will be used, but the functionality will be limited
    and the prescribed slide layouts will have to match either a NxM grid of images (with or without a supplementary
    image such as a legend and with or without a small amount of text) or a single section of text. See
    'allowed_features' below for the allowed attributes, where the 'style' can be 'text' or a 1x1 to 3x3 grid as defined
    in 'image_sets'. Note that images will be scaled to fit on the slide which will likely impact aspect ratios, and if
    an image slide with text is input then the images may be placed over the text.

    If a template is provided then each input slide 'style' must match with the layout name of the appropriate slide
    within the template's Slide Master. For example, if a slide layout was constructed to hold a title, 3 images and a
    section of text with the layout name set as '3x1 with text' then an appropriate dict for the slide would be:
    {'style': '3x1 with text', 'title': 'The title of this slide is...', 'img1': 'firstpic.png', 'img2': 'secondpic.jpg,
    'img3': 'thirdpic.pdf', 'text1': 'Words words words, etc etc etc.'}
    Note that the images (and text if multiple sets of text are defined) will fill the appropriate placeholders from
    left to right, top to bottom.
    '''
    template_provided = False
    allowed_features = ['style', 'title', 'legend', 'img1', 'img2', 'img3', 'img4', 'img5', 'img6', 'img7', 'img8',
                        'img9', 'text1']
    image_sets = {'1x1': [1, 1], '1x2': [1, 2], '1x3': [1, 3], '2x1': [2, 1], '2x2': [2, 2], '2x3': [2, 3],
                  '3x1': [3, 1], '3x2': [3, 2], '3x3': [3, 3]}
    if filename is None:
        filename = 'output_file.pptx'
        print("A filename was not specified. As such, the output will be saved as: %s" %filename)
    if template is None:
        prs = pp.Presentation()
        print("A presentation template was not specified. As such, an automatic template is being used.")
    else:
        if os.path.exists(template):
            prs = pp.Presentation(template)
            template_provided = True
        else:
            prs = pp.Presentation()
            print("The specified template file could not be found. As such, an automatic template is being used.")
    if slides is None or slides == []:
        prs.save(filename)
        print("No slide data was provided. As such, the blank template has been saved as: %s" %filename)
    else:
        if template_provided:
            for s, slide in enumerate(slides):
                for entry in slide.keys():
                    entry.lower()
                num_features = len(slide)
                if num_features > 0:
                    prs = update_custom(prs, slide, s, image_path=image_path)
        else:
            for slide in slides:
                for attr in slide.keys():
                    if attr not in allowed_features:
                        del slide[attr]
                        print("Slide contained an invalid attribute: %s, which has been deleted." %attr)
                num_features = len(slide)
                if num_features > 0:
                    try:
                        format = slide['style']
                    except KeyError:
                        prs = update_fail(prs)
                    if format in image_sets:
                        prs = update_image(prs, slide, arrange=image_sets[format], image_path=image_path)
                    elif format == 'text':
                        prs = update_text(prs, slide)
                    else:
                        prs = update_fail(prs)
                else:
                    prs = update_fail(prs)
        prs.save(filename)
        print("The powerpoint has been saved as: %s" %filename)


def update_image(presentation, slide_details, arrange=[1, 1], image_path=None):
    num_attributes = len(slide_details) - 1
    num_image = arrange[0]*arrange[1]
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    text_present = False
    for entry in slide_details.keys():
        if 'text' in entry or 'par' in entry or 'txt' in entry:
            text_present = True
            break
    if 'title' in slide_details.keys():
        num_attributes -= 1
        title_placeholder = slide.shapes.title
        title_placeholder.text = slide_details['title']
        del slide_details['title']
        title_present = True
    else:
        title_present = False
    if 'legend' in slide_details.keys():
        num_attributes -= 1
        if title_present:
            if text_present:
                left, top, height, width = getproperties(ncols=arrange[0] + 1, nrows=arrange[1], minver=8)
            else:
                left, top, height, width = getproperties(ncols=arrange[0] + 1, nrows=arrange[1], minver=4)
        else:
            if text_present:
                left, top, height, width = getproperties(ncols=arrange[0] + 1, nrows=arrange[1], minver=8)
            else:
                left, top, height, width = getproperties(ncols=arrange[0] + 1, nrows=arrange[1])
        image_name = slide_details['legend']
        img_path = os.path.join(image_path, image_name)
        pic = slide.shapes.add_picture(img_path, left[-1], top[-1], height=height)
        del slide_details['legend']
        legend_present = True
    else:
        if title_present:
            if text_present:
                left, top, height, width = getproperties(ncols=arrange[0], nrows=arrange[1], minver=8)
            else:
                left, top, height, width = getproperties(ncols=arrange[0], nrows=arrange[1], minver=4)
        else:
            if text_present:
                left, top, height, width = getproperties(ncols=arrange[0], nrows=arrange[1], minver=8)
            else:
                left, top, height, width = getproperties(ncols=arrange[0], nrows=arrange[1])
        legend_present = False
    counter = 0
    for image in list(range(num_image)):
        num_attributes -= 1
        image_key = 'img' + str(image + 1)
        if image_key in slide_details.keys():
            if legend_present and (image + 1) % (arrange[0] + 1) == 0:
                counter += 1
            image_name = slide_details[image_key]
            img_path = os.path.join(image_path, image_name)
            pic = slide.shapes.add_picture(img_path, left[counter], top[counter], height=height, width=width)
            del slide_details[image_key]
            counter += 1
        else:
            print("%s image(s) expected in the slide, but %s was either keyed incorrectly or does not exist."
                  %(num_image, image_key))
    if num_attributes > 0:
        presentation = update_text(presentation, slide_details, slide=slide)
    return presentation

def update_text(presentation, slide_details, slide=None):
    num_attributes = len(slide_details) - 1
    if slide is None:
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
    if 'title' in slide_details.keys():
        num_attributes -= 1
        title_placeholder = slide.shapes.title
        title_placeholder.text = slide_details['title']
    if num_attributes > 1:
        for text in list(range(num_attributes)):
            text_key = 'text' + str(text + 1)
            if text_key in slide_details.keys():
                text_info = slide_details[text_key]
                text_frame = slide.placeholders[1].text_frame
                if text == 0:
                    text_frame.clear()
                    p = text_frame.paragraphs[text]
                    run = p.add_run()
                    run.text = text_info
                else:
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = text_info
            else:
                print("%s paragraphs were expected in the slide, but '%s' was either keyed incorrectly or does not "
                      "exist." % (num_attributes, text_key))
    elif num_attributes > 0:
        text_key = 'text1'
        if text_key in slide_details.keys():
            text_info = slide_details[text_key]
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text_info
        else:
            print("A paragraph was expected in the slide, but 'text1' was either keyed incorrectly or does not exist.")
    else:
        print("At least one piece of text was expected for the slide, but none were found.")
    return presentation

def update_fail(presentation):
    slide_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(slide_layout)
    print("The slide format provided was not understood. A blank slide was added in its place. Please provide slide"
          " attributes according to the specifications of savepptx.")
    return presentation

def update_custom(presentation, slide_details, slide_num, image_path=None):
    name = slide_details['style']
    slide = False
    for trial_slide in presentation.slides:
        if trial_slide.slide_layout.name == name:
            slide = trial_slide
    if slide:
        for shape in slide.placeholders:
            if 'title' in slide_details.keys() and 'Title' in shape.name:
                shape.text = slide_details['title']
            if 'Picture' in shape.name or 'Content' in shape.name:
                for entry in slide_details.keys():
                    if 'im' in entry or 'pic' in entry or 'leg' in entry:
                        image_name = slide_details[entry]
                        img_path = os.path.join(image_path, image_name)
                        pic = slide.shapes.add_picture(img_path, shape.left, shape.top, height=shape.height,
                                                       width=shape.width)
                        del slide_details[entry]
                        break
            elif 'Text' in shape.name or 'Content' in shape.name:
                for entry in slide_details.keys():
                    if 'text' in entry or 'par' in entry or 'txt' in entry:
                        shape.text = slide_details[entry]
                        del slide_details[entry]
                        break
        if len(slide_details) > 1:
            for entry in slide_details.keys():
                if 'im' in entry or 'pic' in entry or 'leg' in entry:
                    print('Note: More images were provided for slide %d than there were Picture Placeholders available. '
                          'As such, some images may have been omitted.' % slide_num)
                elif 'text' in entry or 'par' in entry or 'txt' in entry:
                    print('Note: More text paragraphs were provided for slide %d than there were Text Placeholders available'
                          '. As such, some text may have been omitted.' % slide_num)
        elif len(slide_details) > 0:
            if 'im' in slide_details.keys()[0] or 'pic' in slide_details.keys()[0] or 'leg' in slide_details.keys()[0]:
                print('Note: More images were provided for slide %d than there were Picture Placeholders available. '
                      'As such, some images may have been omitted.' % slide_num)
            elif 'text' in slide_details.keys()[0] or 'par' in slide_details.keys()[0] or 'txt' in slide_details.keys()[0]:
                print(
                    'Note: More text paragraphs were provided for slide %d than there were Text Placeholders available'
                    '. As such, some text may have been omitted.' % slide_num)
        return presentation
    else:
        print('Note: The style (%s) attributed to slide %d does not match any layout name in the Slide Master'
              % (name, slide_num))
        return presentation


def getproperties(ncols, nrows, minver=2):
    import numpy
    # image settings for slides
    aspectratio = 1.75
    minhor = 1.0
    availablewidth = 25
    availableheight = 19
    minpercentgap = 0.05

    # image placement calculations
    figheight = float((1 - (nrows - 1) * minpercentgap) * availableheight - minver) / nrows
    figwidth = float((1 - (ncols - 1) * minpercentgap) * availablewidth - minhor) / ncols
    availableaspectratio = figwidth / figheight

    if availableaspectratio > aspectratio:
        ygap = minpercentgap * float(availableheight) / nrows
        figheight = float(availableheight - (nrows - 1) * ygap) / nrows
        figwidth = aspectratio * figheight
        if ncols == 1:
            figwidth = aspectratio * figheight
            xgap = 0
        else:
            xgap = float(availablewidth - ncols * figwidth) / (ncols - 1)
    else:
        xgap = minpercentgap * float(availablewidth) / ncols
        figwidth = float(availablewidth - (ncols - 1) * xgap) / ncols
        figheight = figwidth / float(aspectratio)
        if nrows == 1:
            figheight = figwidth / float(aspectratio)
            ygap = 0
        else:
            ygap = minpercentgap * float(availableheight) / nrows
    while nrows * figheight + (nrows - 1) * ygap + minver > availableheight:
        figheight *= 0.95
        ygap *= 0.95
    while ncols * figwidth + (ncols - 1) * xgap + minhor > availablewidth:
        figwidth *= 0.95
        xgap *= 0.95
    leftcoords = []
    temptop = []
    for j in range(nrows):
        temptop.append(pp.util.Cm(minver + j * (figheight + ygap)))
        for i in range(ncols):
            leftcoords.append(pp.util.Cm(minhor + i * (figwidth + xgap)))

    topcoords = numpy.repeat(temptop, ncols)
    figheight = pp.util.Cm(figheight)
    figwidth = pp.util.Cm(figwidth)

    return leftcoords, topcoords, figheight, figwidth
